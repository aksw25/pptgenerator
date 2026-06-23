import os
import re
import shutil
import subprocess
import tempfile

BAD_PATTERNS = re.compile(
    r"\blorem\b|\bipsum\b|\bplaceholder\b|\bTODO\b|\[insert|\bx{4,}\b",
    re.IGNORECASE,
)


def _text_qa(path: str) -> list[str]:
    issues = []
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        return ["python-pptx not installed — text QA skipped"]

    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        texts.append(run.text)

        combined = " ".join(texts)
        if not combined.strip():
            issues.append(f"Slide {i}: no text content detected")
            continue

        match = BAD_PATTERNS.search(combined)
        if match:
            issues.append(f"Slide {i}: placeholder/lorem text detected near '{match.group()}'")

    return issues


def _visual_qa(path: str) -> list[str]:
    """Convert to images via LibreOffice + pdftoppm, then check with Claude."""
    issues = []

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        return []  # Visual QA unavailable — degrade gracefully

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        return []

    with tempfile.TemporaryDirectory() as tmpdir:
        # Convert pptx → pdf
        result = subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmpdir, path],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            return [f"Visual QA: LibreOffice conversion failed: {result.stderr[:200]}"]

        pdf_name = os.path.splitext(os.path.basename(path))[0] + ".pdf"
        pdf_path = os.path.join(tmpdir, pdf_name)
        if not os.path.exists(pdf_path):
            return ["Visual QA: PDF not produced by LibreOffice"]

        # Convert pdf pages → JPEG
        prefix = os.path.join(tmpdir, "slide")
        result = subprocess.run(
            [pdftoppm, "-jpeg", "-r", "96", pdf_path, prefix],
            capture_output=True, text=True,
        )
        if result.returncode != 0:
            return [f"Visual QA: pdftoppm failed: {result.stderr[:200]}"]

        slides = sorted(f for f in os.listdir(tmpdir) if f.endswith(".jpg"))
        if not slides:
            return ["Visual QA: no slide images produced"]

        issues += _check_slides_with_claude(tmpdir, slides, api_key)

    return issues


def _check_slides_with_claude(tmpdir: str, slides: list[str], api_key: str) -> list[str]:
    import base64
    import anthropic

    client = anthropic.Anthropic(api_key=api_key)
    issues = []

    for fname in slides:
        fpath = os.path.join(tmpdir, fname)
        with open(fpath, "rb") as f:
            b64 = base64.standard_b64encode(f.read()).decode()

        response = client.messages.create(
            model="claude-sonnet-4-6",
            max_tokens=256,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {"type": "base64", "media_type": "image/jpeg", "data": b64},
                        },
                        {
                            "type": "text",
                            "text": (
                                "Inspect this presentation slide image for quality issues. "
                                "Report ONLY problems: text overflowing outside its box, "
                                "overlapping elements, empty/blank areas that should have content, "
                                "or placeholder/lorem text. "
                                "If the slide looks fine, respond with exactly: OK"
                            ),
                        },
                    ],
                }
            ],
        )
        verdict = response.content[0].text.strip()
        if verdict.upper() != "OK":
            issues.append(f"Visual QA {fname}: {verdict}")

    return issues


def run_qa(path: str) -> list[str]:
    issues = _text_qa(path)
    issues += _visual_qa(path)
    return issues
